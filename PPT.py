from pptx import Presentation
from pptx.util import Inches

def createPPT(images: list[str], size: tuple[float, float]=(10, 7.5)):
    """
    default size = (10, 7.5)  4:3
                (13.33, 7.5)  16:9
    """
    
    prs = Presentation()
    canva_w, canva_h = size
    prs.slide_width = Inches(canva_w); prs.slide_height = Inches(canva_h)
    for img in images:
        slide_layout = prs.slide_layouts[5]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.add_picture(img, 0, 0, prs.slide_width, prs.slide_height)
    return prs